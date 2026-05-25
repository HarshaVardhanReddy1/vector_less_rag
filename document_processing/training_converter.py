"""
Training Phase Document Converter
Converts PDF, DOCX, PPTX to clean markdown with proper image handling and formula support

Usage:
    converter = TrainingDocumentConverter()
    converter.convert_file("document.pdf")

Install first:
    Windows:
        winget install pandoc           # for DOCX conversion
        pip install pymupdf pix2tex Pillow python-pptx

    Linux:
        sudo apt-get install -y pandoc
        pip install pymupdf pix2tex Pillow python-pptx

PDF conversion uses PyMuPDF (text + images) + pix2tex (formula regions).
DOCX conversion uses Pandoc.
PPTX conversion uses python-pptx directly.
"""

from pathlib import Path
import subprocess
import io
from pptx import Presentation
from datetime import datetime

import fitz  # PyMuPDF
from PIL import Image
from pix2tex.cli import LatexOCR


# Math fonts used in LaTeX-generated PDFs — presence signals a formula region
_MATH_FONTS = {"CMMI", "CMSY", "CMEX", "SYMBOL", "MT EXTRA", "MATH", "EUCLID"}


class TrainingDocumentConverter:
    """Convert training documents (PDF, DOCX, PPTX) to markdown"""
    
    def __init__(self, output_dir="./markdown_conversions", verbose=True):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        self.verbose = verbose
        self.conversion_log = []
    
    def log(self, message, status="info"):
        """Log conversion progress"""
        if self.verbose:
            emoji = {"info": "ℹ", "success": "✓", "error": "✗", "warning": "⚠"}.get(status, "•")
            print(f"{emoji} {message}")
        
        self.conversion_log.append({
            'timestamp': datetime.now(),
            'message': message,
            'status': status
        })
    
    # ── PDF helpers ───────────────────────────────────────────────────────────

    @staticmethod
    def _block_has_math_font(block: dict) -> bool:
        """Return True if any span in the text block uses a math/formula font."""
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                font_upper = span.get("font", "").upper()
                if any(mf in font_upper for mf in _MATH_FONTS):
                    return True
        return False

    @staticmethod
    def _render_region(page, bbox, zoom: int = 3) -> Image.Image:
        """Render a bounding-box region of a PDF page as a PIL Image."""
        mat = fitz.Matrix(zoom, zoom)
        clip = page.get_pixmap(matrix=mat, clip=fitz.Rect(bbox))
        return Image.open(io.BytesIO(clip.tobytes("png")))

    # ── PDF conversion ────────────────────────────────────────────────────────

    def convert_pdf(self, pdf_path):
        """
        Convert PDF to markdown using PyMuPDF + pix2tex

        Strategy:
        - PyMuPDF extracts text blocks with per-span font metadata
        - Blocks whose fonts match known math fonts (CMMI, CMSY, CMEX …)
          are rendered as images and passed to pix2tex → LaTeX strings
        - Regular text blocks are written as plain markdown
        - Embedded images are saved to an _images/ folder and linked

        Supports:
        - Text extraction (all PDF types)
        - Formula / equation regions → LaTeX ($$...$$)
        - Embedded image extraction
        - No GPU required — pix2tex runs on CPU

        Args:
            pdf_path: Path to PDF file

        Returns:
            Path to generated markdown file, or None if failed
        """

        pdf_path = Path(pdf_path)

        if not pdf_path.exists():
            self.log(f"PDF not found: {pdf_path}", "error")
            return None

        output_md = self.output_dir / f"{pdf_path.stem}.md"
        media_dir = self.output_dir / f"{pdf_path.stem}_images"

        self.log(f"Converting PDF: {pdf_path.name}")
        self.log("Loading LaTeX OCR model (downloads ~200 MB on first run)…")

        try:
            latex_model = LatexOCR()
        except Exception as e:
            self.log(f"Failed to load pix2tex model: {e}", "error")
            self.log("Install with: pip install pix2tex", "error")
            return None

        try:
            doc = fitz.open(str(pdf_path))
        except Exception as e:
            self.log(f"Could not open PDF: {e}", "error")
            return None

        media_dir.mkdir(exist_ok=True)

        md_lines = [f"# {pdf_path.stem}\n\n"]
        image_count = 0
        formula_count = 0

        for page_num, page in enumerate(doc, 1):
            self.log(f"  Processing page {page_num}/{len(doc)}")
            md_lines.append(f"## Page {page_num}\n\n")

            blocks = page.get_text("rawdict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]

            for block in blocks:

                # ── Text block ────────────────────────────────────────────────
                if block["type"] == 0:

                    if self._block_has_math_font(block):
                        # Formula region → render → LaTeX OCR
                        try:
                            img = self._render_region(page, block["bbox"])
                            latex = latex_model(img)
                            md_lines.append(f"$$\n{latex}\n$$\n\n")
                            formula_count += 1
                        except Exception as e:
                            # Fallback: write raw text in code span
                            raw = " ".join(
                                span["text"]
                                for line in block["lines"]
                                for span in line["spans"]
                            ).strip()
                            md_lines.append(f"`{raw}`\n\n")
                            self.log(f"Formula OCR failed on page {page_num}, used raw text: {e}", "warning")

                    else:
                        # Regular text block
                        text = " ".join(
                            span["text"]
                            for line in block["lines"]
                            for span in line["spans"]
                        ).strip()
                        if text:
                            md_lines.append(f"{text}\n\n")

                # ── Embedded image block ──────────────────────────────────────
                elif block["type"] == 1:
                    try:
                        xref = block["image"]
                        base_image = doc.extract_image(xref)
                        img_bytes = base_image["image"]
                        img_ext = base_image["ext"]

                        img_filename = f"page_{page_num:02d}_img_{image_count}.{img_ext}"
                        img_path = media_dir / img_filename
                        with open(img_path, "wb") as f:
                            f.write(img_bytes)

                        md_lines.append(
                            f"![Image page {page_num}]({media_dir.name}/{img_filename})\n\n"
                        )
                        image_count += 1
                    except Exception as e:
                        self.log(f"Could not extract image on page {page_num}: {e}", "warning")

            md_lines.append("---\n\n")

        # Write markdown file
        with open(output_md, "w", encoding="utf-8") as f:
            f.writelines(md_lines)

        self.log(
            f"PDF converted successfully: {output_md.name} "
            f"({formula_count} formulas, {image_count} images extracted)",
            "success"
        )
        return output_md
    
    def convert_docx(self, docx_path):
        """
        Convert DOCX to markdown using Pandoc
        
        Supports:
        - Headers and hierarchy
        - Tables (perfect preservation)
        - Lists and formatting
        - Images with references
        - Formulas from Word
        - Footnotes
        
        Args:
            docx_path: Path to DOCX file
            
        Returns:
            Path to generated markdown file, or None if failed
        """
        
        docx_path = Path(docx_path)
        
        if not docx_path.exists():
            self.log(f"DOCX not found: {docx_path}", "error")
            return None
        
        output_md = self.output_dir / f"{docx_path.stem}.md"
        media_dir = self.output_dir / f"{docx_path.stem}_images"
        
        # Create media directory
        media_dir.mkdir(exist_ok=True)
        
        self.log(f"Converting DOCX: {docx_path.name}")
        
        # Pandoc command - DOCX has excellent support
        cmd = [
            'pandoc',
            str(docx_path),
            '-o', str(output_md),
            '--from=docx',
            '--to=markdown',
            f'--extract-media={media_dir}',
            '--mathjax',  # Preserve formulas
            '--wrap=none',
            '--standalone'
        ]
        
        try:
            result = subprocess.run(cmd, check=True, capture_output=True, text=True)
            self.log(f"DOCX converted successfully: {output_md.name}", "success")
            return output_md
        
        except subprocess.CalledProcessError as e:
            self.log(f"DOCX conversion failed: {e.stderr}", "error")
            return None
        
        except FileNotFoundError:
            self.log("Pandoc not installed. Run: sudo apt-get install -y pandoc", "error")
            return None
    
    def convert_pptx(self, pptx_path):
        """
        Convert PPTX to markdown
        
        Supports:
        - All slides with slide numbers
        - Text from all shapes
        - Image extraction
        - Proper formatting
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            Path to generated markdown file, or None if failed
        """
        
        pptx_path = Path(pptx_path)
        
        if not pptx_path.exists():
            self.log(f"PPTX not found: {pptx_path}", "error")
            return None
        
        output_md = self.output_dir / f"{pptx_path.stem}.md"
        media_dir = self.output_dir / f"{pptx_path.stem}_images"
        
        # Create media directory
        media_dir.mkdir(exist_ok=True)
        
        self.log(f"Converting PPTX: {pptx_path.name}")
        
        try:
            prs = Presentation(pptx_path)
            
            # Start markdown with title
            md_content = f"# {pptx_path.stem}\n\n"
            md_content += f"*PowerPoint Presentation*\n"
            md_content += f"*Total Slides: {len(prs.slides)}*\n\n"
            md_content += "---\n\n"
            
            image_count = 0
            
            # Process each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                md_content += f"## Slide {slide_num}\n\n"
                
                # Extract text from shapes
                text_extracted = False
                for shape in slide.shapes:
                    # Text extraction
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        md_content += f"{text}\n\n"
                        text_extracted = True
                    
                    # Image extraction
                    if shape.shape_type == 13:  # Picture shape type
                        try:
                            image = shape.image
                            image_ext = image.ext
                            image_filename = f"slide_{slide_num:02d}_image_{image_count}.{image_ext}"
                            image_path = media_dir / image_filename
                            
                            # Save image
                            with open(image_path, 'wb') as f:
                                f.write(image.blob)
                            
                            # Add markdown reference
                            image_dir_name = media_dir.name
                            md_content += f"![Slide {slide_num} Image]({image_dir_name}/{image_filename})\n\n"
                            image_count += 1
                            
                        except Exception as e:
                            self.log(f"Could not extract image from slide {slide_num}: {e}", "warning")
                
                if not text_extracted and image_count == 0:
                    md_content += "*[Empty slide]*\n\n"
                
                # Slide separator
                md_content += "---\n\n"
            
            # Write markdown file
            with open(output_md, 'w', encoding='utf-8') as f:
                f.write(md_content)
            
            self.log(f"PPTX converted successfully: {output_md.name} ({image_count} images extracted)", "success")
            return output_md
        
        except Exception as e:
            self.log(f"PPTX conversion failed: {e}", "error")
            return None
    
    def convert_file(self, file_path):
        """
        Auto-detect document type and convert
        
        Args:
            file_path: Path to document file (PDF, DOCX, or PPTX)
            
        Returns:
            Path to generated markdown file, or None if failed
        """
        
        file_path = Path(file_path)
        ext = file_path.suffix.lower()
        
        if ext == '.pdf':
            return self.convert_pdf(file_path)
        
        elif ext == '.docx':
            return self.convert_docx(file_path)
        
        elif ext == '.pptx':
            return self.convert_pptx(file_path)
        
        elif ext == '.doc':
            self.log("Old .doc format detected. Convert to .docx first for better results.", "warning")
            return None
        
        else:
            self.log(f"Unsupported format: {ext}. Supported: .pdf, .docx, .pptx", "error")
            return None
    
    def convert_batch(self, directory, recursive=False):
        """
        Convert all documents in a directory
        
        Args:
            directory: Path to directory containing documents
            recursive: If True, search subdirectories
            
        Returns:
            Dictionary with results {filename: result_path}
        """
        
        directory = Path(directory)
        
        if not directory.exists():
            self.log(f"Directory not found: {directory}", "error")
            return {}
        
        self.log(f"Batch processing directory: {directory}")
        
        results = {}
        pattern = "**/*" if recursive else "*"
        
        for file_path in directory.glob(pattern):
            if file_path.suffix.lower() in ['.pdf', '.docx', '.pptx']:
                self.log(f"Processing: {file_path.name}")
                result = self.convert_file(file_path)
                results[file_path.name] = result
        
        self.log(f"Batch processing complete: {len(results)} files processed", "success")
        return results
    
    def print_summary(self):
        """Print conversion summary"""
        
        success_count = sum(1 for log in self.conversion_log if log['status'] == 'success')
        error_count = sum(1 for log in self.conversion_log if log['status'] == 'error')
        warning_count = sum(1 for log in self.conversion_log if log['status'] == 'warning')
        
        print("\n" + "="*50)
        print("CONVERSION SUMMARY")
        print("="*50)
        print(f"✓ Successful: {success_count}")
        print(f"✗ Errors: {error_count}")
        print(f"⚠ Warnings: {warning_count}")
        print(f"Output directory: {self.output_dir}")
        print("="*50 + "\n")


# Quick usage examples
if __name__ == "__main__":
    
    # Example 1: Convert single PDF
    print("Example 1: Convert single PDF")
    converter = TrainingDocumentConverter()
    result = converter.convert_file("document_processing/docs/comprehensive_test_document.pdf")
    if result:
        print(f"Saved to: {result}\n")
    
    # # Example 2: Convert DOCX
    # print("Example 2: Convert DOCX")
    # result = converter.convert_file("notes.docx")
    # if result:
    #     print(f"Saved to: {result}\n")
    
    # # Example 3: Convert PPTX
    # print("Example 3: Convert PPTX")
    # result = converter.convert_file("presentation.pptx")
    # if result:
    #     print(f"Saved to: {result}\n")
    
    # # Example 4: Batch convert directory
    # print("Example 4: Batch convert directory")
    # results = converter.convert_batch("./training_materials")
    # converter.print_summary()
